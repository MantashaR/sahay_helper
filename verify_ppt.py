import io
from pptx import Presentation
p = Presentation(r'C:/Users/Lenovo/sahay-prototype/Sahay_Quantcraft.pptx')
buf = io.StringIO()
for i, s in enumerate(p.slides, 1):
    buf.write(f"\n=== Slide {i} ===\n")
    for j, shp in enumerate(s.shapes):
        if shp.has_text_frame:
            for k, para in enumerate(shp.text_frame.paragraphs):
                t = "".join(r.text for r in para.runs) or para.text
                if t.strip():
                    buf.write(f"  [{j}] p{k}: {t}\n")
with open(r'C:/Users/Lenovo/sahay-prototype/verify_out.txt', 'w', encoding='utf-8') as f:
    f.write(buf.getvalue())
print("wrote verify_out.txt")
