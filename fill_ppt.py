"""Fill the Quantcraft hackathon PPT template with Sahay content,
preserving the template's fonts/colors/styling by replacing text in
the FIRST run of each existing paragraph and clearing the rest.
"""
import copy
from pptx import Presentation
from pptx.util import Pt
from lxml import etree

SRC = r"C:/Users/Lenovo/Downloads/6a0613a644f7c_Quantcraft_PPT_template.pptx"
DST = r"C:/Users/Lenovo/sahay-prototype/Sahay_Quantcraft.pptx"

# ---------------- content ----------------
# Each entry: shape_index_on_slide -> list of paragraph strings.
# Paragraph count should match the template's existing paragraph count
# so styling carries through. Extra paragraphs are appended cloned from the last.

SLIDES = {
    1: {  # Title slide
        4: ["Team"],
        5: [
            "Team Sahay",
            "Member 1  —  Role",
            "Member 2  —  Role",
            "Member 3  —  Role",
            "Member 4  —  Role",
        ],
    },
    2: {  # Problem statement & theme
        0: ["Problem Statement & Theme"],
        1: [
            "Problem:  India runs 3,000+ welfare schemes worth ₹15+ lakh crore — yet under 30% of eligible citizens ever claim them. Discovery is broken: portals are English-first, jargon-heavy, and scattered across 50+ ministry websites.",
            "Impact:  ~800 million Bharat citizens — farmers, widows, gig workers, artisans, the elderly — leave entitled benefits on the table every year. A widow in UP doesn't know a ₹6,000/yr pension exists in her name.",
            "Theme Alignment:  AI for Social Good × Bharat-first. A multilingual, voice-enabled AI navigator that turns natural Hindi/Hinglish speech into a ranked list of schemes the user actually qualifies for.",
        ],
    },
    3: {  # Solution
        6: ["Solution"],
        7: [
            "",
            "Solution Overview:  Sahay (सहाय) is an AI Welfare Navigator. The user types or speaks freely in Hindi / Hinglish / English — \"58 saal, kisan, 1.5 acre, beti ki padhai\" — and Sahay returns the top schemes with eligibility, documents required, and how-to-apply links.",
            "Key Features:  (1) Bilingual chat + voice input  (2) Signal-extraction engine that catches age, occupation, land-size, life-events from free text  (3) 25 real central-govt schemes (PM-KISAN, Ayushman Bharat, MUDRA, Ujjwala, PMAY, Vishwakarma, SSY, NSAP pensions, e-Shram…)  (4) One-click demo personas — Ramesh, Lakshmi, Priya, Irfan  (5) Animated total-₹-unlocked counter that quantifies the impact live.",
            "Unique Value:  Most schemes-search tools are keyword forms in English. Sahay is conversational, vernacular-first, and returns a *reasoned* match — every card explains *why* this scheme fits this person.",
        ],
    },
    4: {  # Architecture
        6: ["Architecture"],
        7: [
            "Tech Stack:  Python + Flask backend  •  Vanilla JS / HTML / CSS frontend  •  Claude Haiku 4.5 for AI enrichment  •  JSON scheme corpus (offline-first, zero-DB).",
            "Data Flow:  User speaks/types in Hindi or English → signal extractor pulls age, occupation, land, life-events → scheme scorer ranks 25 schemes → build_why() attaches plain-language reasons → frontend renders cards + ₹-counter.",
            "APIs:  Anthropic Claude API (Haiku 4.5) for \"why this scheme\" copy  •  Web Speech API for voice input  •  Internal /api/chat, /api/schemes, /api/personas.",
        ],
    },
    5: {  # Market research
        6: ["Market Research & Analysis"],
        7: [
            "Target Users:  ~800M citizens of Bharat — farmers, widows, senior citizens, gig & platform workers, MSME artisans, students from low-income families. Primary entry point: rural & tier-2/3 smartphones via WhatsApp-style chat or voice.",
            "Market Size:  3,000+ central + state schemes  •  ₹15+ lakh crore total annual welfare outlay  •  ~70% leakage / under-claim gap = ₹10+ lakh crore worth of unclaimed benefits sitting on the table every year.",
            "Revenue Model:  (1) Govt / state-mission licensing (white-label for Common Service Centres, Panchayats)  (2) CSR + foundation grants tied to claim-conversion metrics  (3) NBFC / fintech partnerships — qualified-lead handoff for MUDRA, Stand-Up India, KCC loans.",
            "Expected Growth:  Pilot in 1 state via CSC network → 10M users in 12 months → scale to all 22 official languages and 250K Common Service Centres → measurable lift in scheme-claim rate is the north-star metric.",
        ],
    },
    6: {  # Future scope
        6: ["Future Scope"],
        7: [
            "WhatsApp & IVR Channels:  Meet Bharat where it already is — Sahay over WhatsApp Business API + a toll-free IVR for feature-phone users with no smartphone.",
            "Full 22-Language Coverage:  Extend signal-map and scheme copy to all 8th-Schedule languages — Marathi, Tamil, Telugu, Bengali, Punjabi, Odia, Assamese, Kannada, Malayalam, Gujarati, and more.",
            "Document-Auto-Fill & Application Tracking:  OCR Aadhaar/ration card → pre-fill scheme forms → submit to UMANG/state portals → track application status end-to-end inside Sahay.",
            "State-Scheme + Private-Welfare Expansion:  Grow from 25 central schemes to 3,000+ central + state + corporate-CSR schemes, with eligibility rules versioned per state.",
            "Impact Dashboards for Govt Partners:  Live dashboards for Ministries / NITI Aayog showing district-level claim-rate uplift, scheme demand heatmaps, and unmet-need signals.",
        ],
    },
    7: {  # TRAE Account
        6: ["TRAE Account Upholding"],
        7: [
            "Screenshot of TRAE dashboard for each teammate — to be added before submission.",
            "Username visible in every screenshot.",
            "Proof of participation for all teammates attached on this slide.",
            "Placeholder: insert screenshots in the image area below.",
        ],
    },
    8: {  # Thank you — leave as is
    },
}


def replace_paragraph_text(p_elem, new_text, template_run=None):
    """Replace all text in a paragraph element with new_text, keeping the
    formatting of the first existing run (or a provided template run)."""
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    runs = p_elem.findall('a:r', nsmap)
    # find the run to keep as the styling template
    keep = template_run if template_run is not None else (runs[0] if runs else None)
    # remove every existing run
    for r in runs:
        p_elem.remove(r)
    # also strip any field runs (a:fld) so they don't leave stray text
    for fld in p_elem.findall('a:fld', nsmap):
        p_elem.remove(fld)
    if keep is None:
        # build a minimal run if the paragraph had no runs at all
        r = etree.SubElement(p_elem, '{http://schemas.openxmlformats.org/drawingml/2006/main}r')
        t = etree.SubElement(r, '{http://schemas.openxmlformats.org/drawingml/2006/main}t')
        t.text = new_text
        return
    new_run = copy.deepcopy(keep)
    t = new_run.find('a:t', nsmap)
    if t is None:
        t = etree.SubElement(new_run, '{http://schemas.openxmlformats.org/drawingml/2006/main}t')
    t.text = new_text
    p_elem.append(new_run)


def fill_shape(shape, paragraphs):
    """Fill a text-frame shape with the given list of paragraph strings,
    preserving existing paragraph styles. Adds cloned paragraphs if needed."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    txBody = tf._txBody
    existing = txBody.findall('a:p', nsmap)
    # template paragraph to clone when we need extra rows
    template_p = existing[-1] if existing else None

    # write into existing paragraphs
    for i, text in enumerate(paragraphs):
        if i < len(existing):
            replace_paragraph_text(existing[i], text)
        else:
            new_p = copy.deepcopy(template_p)
            # remove all runs/flds from clone first
            for child in list(new_p):
                tag = etree.QName(child).localname
                if tag in ('r', 'fld'):
                    new_p.remove(child)
            # find a run we can clone from the template_p for styling
            tpl_runs = template_p.findall('a:r', nsmap)
            tpl_run = tpl_runs[0] if tpl_runs else None
            if tpl_run is not None:
                new_run = copy.deepcopy(tpl_run)
                t = new_run.find('a:t', nsmap)
                if t is None:
                    t = etree.SubElement(new_run, '{http://schemas.openxmlformats.org/drawingml/2006/main}t')
                t.text = text
                new_p.append(new_run)
            else:
                r = etree.SubElement(new_p, '{http://schemas.openxmlformats.org/drawingml/2006/main}r')
                t = etree.SubElement(r, '{http://schemas.openxmlformats.org/drawingml/2006/main}t')
                t.text = text
            txBody.append(new_p)

    # if template had MORE paragraphs than we want, blank the extras
    for j in range(len(paragraphs), len(existing)):
        replace_paragraph_text(existing[j], "")


def main():
    p = Presentation(SRC)
    for slide_idx, shape_map in SLIDES.items():
        slide = p.slides[slide_idx - 1]
        shapes = list(slide.shapes)
        for shape_idx, paragraphs in shape_map.items():
            fill_shape(shapes[shape_idx], paragraphs)
    p.save(DST)
    print(f"Saved -> {DST}")


if __name__ == "__main__":
    main()
