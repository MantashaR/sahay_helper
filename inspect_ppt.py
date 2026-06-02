from pptx import Presentation
from pptx.util import Emu

p = Presentation(r'C:/Users/Lenovo/Downloads/6a0613a644f7c_Quantcraft_PPT_template.pptx')
print(f"Slide size: {p.slide_width} x {p.slide_height}  ({Emu(p.slide_width).inches:.2f} x {Emu(p.slide_height).inches:.2f} in)")
for i, s in enumerate(p.slides, 1):
    print(f"\n=== Slide {i}  layout='{s.slide_layout.name}' ===")
    for j, shp in enumerate(s.shapes):
        info = f"  [{j}] {shp.shape_type} name='{shp.name}'"
        try:
            info += f"  pos=({Emu(shp.left).inches:.2f},{Emu(shp.top).inches:.2f}) size=({Emu(shp.width).inches:.2f}x{Emu(shp.height).inches:.2f})"
        except Exception:
            pass
        print(info)
        if shp.has_text_frame:
            for k, para in enumerate(shp.text_frame.paragraphs):
                txt = "".join(r.text for r in para.runs) or para.text
                if txt.strip():
                    print(f"      p{k}: {txt!r}")
